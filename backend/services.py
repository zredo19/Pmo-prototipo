"""
Core business logic for PMO Agent.
Handles file parsing, AI analysis, and priority calculations.
"""

import json
import re
import os
from io import BytesIO
from typing import Any

import pandas as pd
from groq import AsyncGroq
from pptx import Presentation
from pptx.exc import PackageNotFoundError

from database import load_env_file
load_env_file()


class FileParsingError(Exception):
    """Raised when file parsing fails."""
    pass


class AIAnalysisError(Exception):
    """Raised when AI analysis fails."""
    pass


def parse_excel(file_content: bytes) -> dict[str, Any]:
    """
    Parse Excel file and extract data summary.
    Handles NaN values gracefully by converting to None.
    
    Args:
        file_content: Raw bytes of the Excel file
        
    Returns:
        Dictionary with sheet names as keys and data summaries
        
    Raises:
        FileParsingError: If file cannot be parsed
    """
    try:
        excel_file = BytesIO(file_content)
        excel_data = pd.ExcelFile(excel_file)
        
        result = {
            "sheets": [],
            "total_rows": 0,
            "summary": ""
        }
        
        all_data_summary = []
        
        for sheet_name in excel_data.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            
            # Replace NaN with None for JSON serialization
            df = df.where(pd.notnull(df), None)
            
            sheet_info = {
                "name": sheet_name,
                "rows": len(df),
                "columns": list(df.columns),
                "preview": df.head(10).to_dict(orient="records"),
            }
            
            result["sheets"].append(sheet_info)
            result["total_rows"] += len(df)
            
            # Create text summary for AI
            summary_text = f"\n=== Sheet: {sheet_name} ===\n"
            summary_text += f"Columns: {', '.join(map(str, df.columns))}\n"
            summary_text += f"Total rows: {len(df)}\n"
            
            # Include numeric columns statistics
            numeric_cols = df.select_dtypes(include=["number"]).columns
            for col in numeric_cols:
                summary_text += f"{col}: min={df[col].min()}, max={df[col].max()}, mean={df[col].mean():.2f}\n"
            
            # Include first few rows as context
            summary_text += "\nFirst 5 rows:\n"
            summary_text += df.head(5).to_string()
            
            all_data_summary.append(summary_text)
        
        result["summary"] = "\n".join(all_data_summary)
        
        result["summary"] = "\n".join(all_data_summary)
        # Store raw dataframe for intelligent filtering later if needed
        # We'll attach it to the first sheet for simplicity in this single-sheet context
        if not excel_data.sheet_names:
             result["df"] = pd.DataFrame()
        else:
             result["df"] = pd.read_excel(excel_file, sheet_name=excel_data.sheet_names[0]) # Start with first sheet
        
        return result
        
    except Exception as e:
        raise FileParsingError(f"Failed to parse Excel file: {str(e)}")


def extract_pptx_text(file_content: bytes) -> dict[str, Any]:
    """
    Extract all text content from PowerPoint presentation.
    
    Args:
        file_content: Raw bytes of the PPTX file
        
    Returns:
        Dictionary with slide texts and summary
        
    Raises:
        FileParsingError: If file cannot be parsed or is empty
    """
    try:
        pptx_file = BytesIO(file_content)
        presentation = Presentation(pptx_file)
        
        slides_text = []
        all_text = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            slide_content = "\n".join(slide_texts)
            slides_text.append({
                "slide_number": slide_num,
                "content": slide_content
            })
            
            if slide_content:
                all_text.append(f"=== Slide {slide_num} ===\n{slide_content}")
        
        if not all_text:
            raise FileParsingError("PowerPoint file is empty or contains no readable text.")
        
        return {
            "slides": slides_text,
            "total_slides": len(presentation.slides),
            "summary": "\n\n".join(all_text)
        }
        
    except PackageNotFoundError:
        raise FileParsingError("Invalid or corrupted PowerPoint file.")
    except FileParsingError:
        raise
    except Exception as e:
        raise FileParsingError(f"Failed to parse PowerPoint file: {str(e)}")


def extract_project_names_from_pptx_text(pptx_text: str) -> list[str]:
    """
    Extract potential project names from PPTX text using simple heuristics to filter Excel data.
    Looking for lines that look like titles or capitalized phrases.
    """
    potential_names = []
    lines = pptx_text.split('\n')
    for line in lines:
        clean_line = line.strip()
        # Heuristic: Project names are usually short (5-50 chars), often start with capital letters
        if 5 < len(clean_line) < 60 and clean_line[0].isupper():
            potential_names.append(clean_line)
    return potential_names


def filter_excel_data_by_relevance(df: pd.DataFrame, pptx_text: str) -> str:
    """
    Filter Excel rows to include only those that might match content in the PPTX.
    This creates a 'context window' for the AI with relevant data only.
    """
    relevant_rows = []
    
    # Simple keyword matching: Check if any part of project name exists in PPT text
    # or if PPT text contains the Project ID
    
    # Convert PPT text to lowercase for case-insensitive matching
    pptx_lower = pptx_text.lower()
    
    for index, row in df.iterrows():
        # Construct search terms from row data
        search_terms = []
        if pd.notna(row.get('Project_Name')):
            search_terms.append(str(row['Project_Name']).lower())
        if pd.notna(row.get('Project_ID')):
            search_terms.append(str(row['Project_ID']).lower())
            
        # Check if row is relevant
        is_relevant = False
        for term in search_terms:
            if term in pptx_lower:
                is_relevant = True
                break
        
        if is_relevant:
            relevant_rows.append(row)
            
    if not relevant_rows:
        return "No exact row matches found based on text search. Providing top 20 rows." + df.head(20).to_string()
        
    # Create DataFrame from relevant rows
    filtered_df = pd.DataFrame(relevant_rows)
    return f"Found {len(filtered_df)} relevant matching rows in Excel:\n" + filtered_df.to_string()



from groq import AsyncGroq

# ... (imports remain the same, removing 'from groq import Groq' if it was there, effectively replacing the block)

async def analyze_with_groq(excel_data_obj: dict[str, Any], pptx_summary: str) -> dict[str, Any]:
    """
    Use Groq AI to analyze discrepancies between Excel data and PowerPoint content.
    Uses intelligent filtering (RAG-lite) to send only relevant Excel rows to the LLM.
    
    Args:
        excel_data_obj: Result dictionary from parse_excel containing 'df' and 'summary'
        pptx_summary: Text summary of PowerPoint content
        
    Returns:
        Dictionary with analysis results including discrepancies list
        
    Raises:
        AIAnalysisError: If AI analysis fails
    """
    
    # 1. Prepare relevant Excel Context using RAG-lite
    if 'df' in excel_data_obj and not excel_data_obj['df'].empty:
        df = excel_data_obj['df']
        excel_context = filter_excel_data_by_relevance(df, pptx_summary)
    else:
        # Fallback to the original summary if dataframe is missing
        excel_context = excel_data_obj.get('summary', 'No Excel data found.')
    api_key = os.getenv("GROQ_API_KEY")
    
    if not api_key or api_key == "your_groq_api_key_here":
        raise AIAnalysisError(
            "GROQ_API_KEY no configurada. Por favor configura una llave válida en el archivo .env."
        )
    
    try:
        client = AsyncGroq(api_key=api_key)
        
        prompt = f"""Eres un asistente analista de datos experto. Analiza los siguientes datos de Excel y contenido de PowerPoint en busca de inconsistencias.

Enfócate en:
1. Números de presupuesto que no coinciden entre documentos
2. Fechas que son inconsistentes
3. Nombres de proyectos o descripciones que difieren
4. Valores numéricos que se contradicen
5. Datos faltantes en un documento que existen en el otro

REGLAS CRÍTICAS PARA LA COMPARACIÓN:
- NO reportes una discrepancia si los valores son idénticos.
- NO reportes una discrepancia si los valores son numéricamente equivalentes (ej: "1744000" vs 1744000).
- NO reportes diferencias de formato (ej: "2024-01-01" vs "01/01/2024" si representan la misma fecha).
- SOLO reporta diferencias sustanciales que requieran atención del usuario.
- RESPONDE SIEMPRE EN ESPAÑOL.

DATOS EXCEL:
DATOS EXCEL (Filtrados por relevancia):
{excel_context}

CONTENIDO POWERPOINT:
{pptx_summary}

Retorna tu análisis ESTRICTAMENTE como un objeto JSON con esta estructura exacta:
{{
    "discrepancies": [
        {{
            "type": "presupuesto|fecha|texto|faltante|numerico",
            "severity": "alta|media|baja",
            "description": "Descripción clara de la discrepancia en Español",
            "excel_value": "Valor encontrado en Excel (o null)",
            "pptx_value": "Valor encontrado en PowerPoint (or null)",
            "recommendation": "Acción sugerida para resolver en Español"
        }}
    ],
    "summary": "Breve resumen general de los hallazgos en Español",
    "match_score": 85
}}

Si no se encuentran discrepancias, retorna un array 'discrepancies' vacío.
IMPORTANTE: Retorna SOLO JSON válido, sin texto adicional ni markdown."""

        response = await client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "Eres un analista de datos preciso. Responde siempre con JSON válido y en Español."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.0,  # Determinístico: mismos archivos = mismos resultados
            max_tokens=2000,
        )
        
        response_text = response.choices[0].message.content.strip()
        
        # Try to parse JSON response
        # Try to parse JSON response using robust regex extraction
        try:
            # Look for the first valid JSON object in the response
            # Pattern: matches { ... } including nested braces, minimally
            # Using simple curly brace matching for robustness
            match = re.search(r'(\{.*\})', response_text, re.DOTALL)
            
            if match:
                json_str = match.group(1)
                result = json.loads(json_str)
            else:
                # Fallback: try raw loads if no braces found (unlikely for valid JSON)
                result = json.loads(response_text)
            
            # Validate required fields
            if "discrepancies" not in result:
                result["discrepancies"] = []
            
            # Post-processing: Filter out hallucinations where values are identical
            # This is a critical "Safety Layer" because LLMs sometimes hallucinate differences
            valid_discrepancies = []
            for disc in result["discrepancies"]:
                val_excel = str(disc.get("excel_value", "")).strip().replace("$", "").replace(",", "")
                val_pptx = str(disc.get("pptx_value", "")).strip().replace("$", "").replace(",", "")
                
                # Only keep if they are substantially different
                if val_excel != val_pptx:
                    valid_discrepancies.append(disc)
            
            result["discrepancies"] = valid_discrepancies

            if "summary" not in result:
                result["summary"] = "Análisis completado."
            if "match_score" not in result:
                result["match_score"] = 100 - (len(result["discrepancies"]) * 10)
            
            return result
            
        except json.JSONDecodeError:
            # If JSON parsing fails, create structured response from text
            return {
                "discrepancies": [{
                    "type": "texto",
                    "severity": "media",
                    "description": response_text[:500],
                    "excel_value": None,
                    "pptx_value": None,
                    "recommendation": "Se recomienda revisión manual"
                }],
                "summary": "Análisis IA completado con notas de análisis.",
                "match_score": 70
            }
            
    except Exception as e:
        raise AIAnalysisError(f"Error de API Groq: {str(e)}")


def process_excel_batch(file_content: bytes) -> dict[str, Any]:
    """
    Parse Excel file and calculate priority scores for all projects.
    
    Args:
        file_content: Raw bytes of the Excel file
        
    Returns:
        Dictionary with list of project results and summary stats
    """
    try:
        excel_file = BytesIO(file_content)
        df = pd.read_excel(excel_file)
        
        # Normalize column names to snake_case
        df.columns = [
            str(col).strip().lower().replace(' ', '_').replace('-', '_') 
            for col in df.columns
        ]
        
        results = []
        
        for _, row in df.iterrows():
            # Extract basic info
            project_id = row.get('id_proyecto', 'Desconocido')
            name = row.get('nombre_proyecto', 'Sin Título')
            
            # Extract metrics with fallbacks
            roi = float(row.get('metric_roi_0_100', 0))
            urgency = float(row.get('metric_urgencia_0_100', 0))
            risk = float(row.get('metric_riesgo_0_100', 50))
            alignment = float(row.get('metric_alineacion_0_100', 50))
            resources = float(row.get('metric_recursos_0_100', 50))
            
            # Calculate Score
            score_data = calculate_priority_score(
                roi=roi,
                urgency=urgency,
                risk=risk,
                strategic_alignment=alignment,
                resource_availability=resources
            )
            
            # Merge project info with score data
            project_result = {
                "id": project_id,
                "name": name,
                "sponsor": row.get('sponsor_ejecutivo', 'N/A'),
                "area": row.get('area_negocio', 'N/A'),
                **score_data
            }
            results.append(project_result)
            
        # Sort by score descending
        results.sort(key=lambda x: x['score'], reverse=True)
        
        return {
            "total_projects": len(results),
            "results": results
        }
        
    except Exception as e:
        raise FileParsingError(f"Error al procesar lote Excel: {str(e)}")



def calculate_priority_score(
    roi: float,
    urgency: float,
    risk: float,
    strategic_alignment: float = 50.0,
    resource_availability: float = 50.0
) -> dict[str, Any]:
    """
    Calculate project priority score using weighted algorithm.
    
    Args:
        roi: Expected Return on Investment (0-100)
        urgency: Time sensitivity of the project (0-100)
        risk: Risk level where higher = riskier (0-100)
        strategic_alignment: Alignment with company strategy (0-100)
        resource_availability: Available resources for project (0-100)
        
    Returns:
        Dictionary with priority score and breakdown
    """
    # Validate inputs
    def clamp(value: float, min_val: float = 0.0, max_val: float = 100.0) -> float:
        return max(min_val, min(max_val, float(value)))
    
    roi = clamp(roi)
    urgency = clamp(urgency)
    risk = clamp(risk)
    strategic_alignment = clamp(strategic_alignment)
    resource_availability = clamp(resource_availability)
    
    # Weights for each factor (total = 1.0)
    weights = {
        "roi": 0.30,
        "urgency": 0.25,
        "risk_adjusted": 0.20,
        "strategic_alignment": 0.15,
        "resource_availability": 0.10,
    }
    
    # Risk is inverted (lower risk = higher score)
    risk_adjusted = 100 - risk
    
    # Calculate weighted score
    weighted_scores = {
        "roi": roi * weights["roi"],
        "urgency": urgency * weights["urgency"],
        "risk_adjusted": risk_adjusted * weights["risk_adjusted"],
        "strategic_alignment": strategic_alignment * weights["strategic_alignment"],
        "resource_availability": resource_availability * weights["resource_availability"],
    }
    
    total_score = sum(weighted_scores.values())
    
    # Determine priority tier
    if total_score >= 80:
        tier = "Crítico"
        tier_color = "#DC2626"  # red-600
    elif total_score >= 60:
        tier = "Alto"
        tier_color = "#EA580C"  # orange-600
    elif total_score >= 40:
        tier = "Medio"
        tier_color = "#CA8A04"  # yellow-600
    else:
        tier = "Bajo"
        tier_color = "#16A34A"  # green-600
    
    return {
        "score": round(total_score, 1),
        "tier": tier,
        "tier_color": tier_color,
        "breakdown": {
            "roi": {"value": roi, "weight": weights["roi"], "contribution": round(weighted_scores["roi"], 1)},
            "urgency": {"value": urgency, "weight": weights["urgency"], "contribution": round(weighted_scores["urgency"], 1)},
            "risk": {"value": risk, "weight": weights["risk_adjusted"], "contribution": round(weighted_scores["risk_adjusted"], 1)},
            "strategic_alignment": {"value": strategic_alignment, "weight": weights["strategic_alignment"], "contribution": round(weighted_scores["strategic_alignment"], 1)},
            "resource_availability": {"value": resource_availability, "weight": weights["resource_availability"], "contribution": round(weighted_scores["resource_availability"], 1)},
        },
        "recommendation": get_priority_recommendation(tier, roi, urgency, risk)
    }


def get_priority_recommendation(tier: str, roi: float, urgency: float, risk: float) -> str:
    """Generate recommendation based on priority analysis."""
    if tier == "Crítico":
        return "Acción inmediata requerida. Asignar recursos y comenzar implementación."
    elif tier == "Alto":
        if roi > 70:
            return "Proyecto de alto ROI. Priorizar para el siguiente sprint o ciclo."
        elif urgency > 70:
            return "Proyecto sensible al tiempo. Acelerar proceso de aprobación."
        else:
            return "Fuerte candidato para ejecución prioritaria. Revisar asignación de recursos."
    elif tier == "Medio":
        if risk > 60:
            return "Prioridad moderada con riesgo elevado. Considerar mitigación de riesgos antes de proceder."
        else:
            return "Encolar para ejecución estándar. Monitorear cambios de prioridad."
    else:
        if roi < 30:
            return "Bajo ROI. Considerar despriorizar o revisar alcance."
        else:
            return "Programar para consideración futura. Documentar para revisión trimestral."
