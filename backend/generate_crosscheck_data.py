# ... (imports)
import pandas as pd
import random
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Setup data
# Setup data
n_rows = 100
project_types = ["Infraestructura", "Software", "Cumplimiento", "Marketing", "RRHH", "I+D"]
managers = ["J. Pérez", "M. García", "A. López", "L. Rodríguez", "K. Muñoz"]
statuses = ["Rojo", "Amarillo", "Verde", "Completado", "En Pausa"]

data = {
    "Project_ID": [f"PROJ-{2024000+i}" for i in range(n_rows)],
    "Project_Name": [],
    "Department": [],
    "Project_Manager": [],
    "Total_Budget_Approved": [],
    "Actual_Spend_YTD": [],
    "Forecast_EAC": [],
    "Status_RAG": [],
    "Completion_Percent": [],
    "Go_Live_Date": [],
    "Comments": []
}

prefixes = ["Global", "Regional", "Local", "Estratégico", "Interno"]
nouns = ["Transformación", "Migración", "Actualización", "Despliegue", "Auditoría", "Iniciativa", "Plataforma", "Panel"]

for i in range(n_rows):
    # Name generation
    p_type = random.choice(project_types)
    name = f"{random.choice(nouns)} {p_type} {random.choice(prefixes)} {random.randint(1,99)}"
    data["Project_Name"].append(name)
    data["Department"].append(p_type)
    data["Project_Manager"].append(random.choice(managers))
    
    # Financials
    budget = round(random.uniform(50000, 2000000), -3) # Round to nearest 1000
    spend = round(budget * random.uniform(0.1, 1.1), -2)
    forecast = round(spend + (budget - spend) * random.uniform(0.9, 1.2), -2)
    
    data["Total_Budget_Approved"].append(budget)
    data["Actual_Spend_YTD"].append(spend)
    data["Forecast_EAC"].append(forecast)
    
    # Status
    # Weighted status to have enough Red/Yellow
    status = random.choices(statuses, weights=[0.2, 0.3, 0.3, 0.1, 0.1], k=1)[0]
    data["Status_RAG"].append(status)
    
    # Completion
    comp = random.randint(0, 100)
    data["Completion_Percent"].append(comp)
    
    # Dates
    date = datetime(2024, 1, 1) + timedelta(days=random.randint(0, 365))
    data["Go_Live_Date"].append(date.strftime("%Y-%m-%d"))
    
    # Comments
    if status == "Rojo":
        data["Comments"].append("Riesgo de sobrecosto detectado.")
    elif status == "Amarillo":
        data["Comments"].append("Restricciones de recursos.")
    else:
        data["Comments"].append("En camino.")

df = pd.DataFrame(data)

# Save Excel
output_excel = "dataset_pmo_crosscheck_100.xlsx"
df.to_excel(output_excel, index=False)
print(f"Excel creado: {output_excel}")

# --- Generate Matching PPTX with Discrepancies ---
prs = Presentation()

# Introduction Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Reporte de Estado de Proyectos Mensual"
subtitle.text = f"Generado: {datetime.now().strftime('%Y-%m-%d')}"

# Select top 10 critical projects (Red/Yellow) to feature
critical_projects = df[df['Status_RAG'].isin(['Rojo', 'Amarillo'])].nlargest(10, 'Total_Budget_Approved')

# Content Slides
bullet_slide_layout = prs.slide_layouts[1]

# Prompt generation lists
prompt_lines = []
prompt_lines.append("Actúa como un Asistente PMO. Genera una presentación de PowerPoint (o el contenido para una) basada en estos datos críticos.")
prompt_lines.append("Instrucciones: Crea diapositivas para los siguientes proyectos. Introduce ALGUNOS errores sutiles en presupuesto o fechas para probar la atención al detalle.")
prompt_lines.append("\nDatos Críticos del Excel (Fuente de Verdad):")

for _, row in critical_projects.iterrows():
    # PPT Generation
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    # Title
    title_shape = shapes.title
    title_shape.text = row['Project_Name']
    
    # Body
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    # Randomly introduce errors matches the prompt request
    # 1. Budget
    pptx_budget = row['Total_Budget_Approved']
    budget_note = ""
    if random.random() < 0.3: # 30% chance error
        factor = random.choice([0.8, 1.2, 1.5]) # Discrepancy
        pptx_budget = round(pptx_budget * factor, -3)
        budget_note = " (ERROR INTENCIONAL)"
        
    p = tf.add_paragraph()
    p.text = f"Presupuesto Aprobado: ${pptx_budget:,.0f}"
    
    # 2. Date
    pptx_date = row['Go_Live_Date']
    date_note = ""
    if random.random() < 0.3: # 30% chance error
        original_date = datetime.strptime(row['Go_Live_Date'], "%Y-%m-%d")
        new_date = original_date + timedelta(days=random.choice([45, -45]))
        pptx_date = new_date.strftime("%Y-%m-%d")
        date_note = " (ERROR INTENCIONAL)"
        
    p = tf.add_paragraph()
    p.text = f"Fecha Go-Live: {pptx_date}"
    
    # 3. Manager
    p = tf.add_paragraph()
    p.text = f"Project Manager: {row['Project_Manager']}"
    
    # 4. Status
    p = tf.add_paragraph()
    p.text = f"Estado: {row['Status_RAG']}"

    # Add to Prompt Text (Simulating what user would ask another AI)
    # We provide the 'source of truth' but tell the AI to be 'realistic' (which implies potential human error simulation if we want, OR we provide the modified values if we want the user to generate THIS specific bad PPT)
    # The user asked for "the prompt to create a presentation ppt... make there be some errors".
    # So the PROMPT ITSELF should probably contain the MODIFIED data if the user uses a "Create PPT from this text" tool.
    # OR the prompt instructs the AI to make errors.
    # Let's give the user a prompt with the *correct* data and an instruction to *add* errors, OR give the prompt with the *already modified* data so the AI just renders it.
    # Giving the modified data is safer for cross-check validation because we know exactly what errors define.
    
    prompt_lines.append(f"\nProyecto: {row['Project_Name']}")
    prompt_lines.append(f"- Presupuesto: ${pptx_budget:,.0f}{budget_note}")
    prompt_lines.append(f"- Fecha: {pptx_date}{date_note}")
    prompt_lines.append(f"- Manager: {row['Project_Manager']}")
    prompt_lines.append(f"- Estado: {row['Status_RAG']}")

# Save PPTX
output_pptx = "reporte_estado_proyectos_100.pptx"
prs.save(output_pptx)
print(f"PowerPoint creado: {output_pptx}")

# Save Prompt
output_prompt = "prompt_para_ia_ppt.txt"
with open(output_prompt, "w", encoding="utf-8") as f:
    f.write("\n".join(prompt_lines))
print(f"Prompt creado: {output_prompt}")
